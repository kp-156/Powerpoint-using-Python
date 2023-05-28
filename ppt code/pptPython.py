from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Load the font file
font_file = 'sample_font_file.ttf'

# Read content for Slide1
with open('slide1_content.txt', 'r') as f:
    slide1_content = f.read()

# Read content for Slide2
with open('sample_slide2_input.txt', 'r') as f:
    slide2_content = f.read()

# Create a presentation object
presentation = Presentation()

# Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

# Set the font for Slide 1
title.text_frame.clear()
title.text_frame.add_paragraph().text = 'PPT using Python'
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title.text_frame.paragraphs[0].font.name = font_file
title.text_frame.paragraphs[0].font.size = Pt(32)

subtitle.text_frame.clear()
subtitle.text_frame.add_paragraph().text = slide1_content
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
subtitle.text_frame.paragraphs[0].font.name = font_file
subtitle.text_frame.paragraphs[0].font.size = Pt(24)

# Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide2.shapes.title
content = slide2.placeholders[1]
 
# Set the font for Slide 2
title.text_frame.clear()
title.text_frame.add_paragraph().text = 'Slide 2'
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title.text_frame.paragraphs[0].font.name = font_file
title.text_frame.paragraphs[0].font.size = Pt(32)

content.text_frame.clear()
content.text_frame.add_paragraph().text = slide2_content
content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
content.text_frame.paragraphs[0].font.name = font_file
content.text_frame.paragraphs[0].font.size = Pt(24)

# Save the presentation
presentation.save('output.pptx')
